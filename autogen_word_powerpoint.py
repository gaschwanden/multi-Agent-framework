from docx import Document  # For Word document creation
from pptx import Presentation  # For PowerPoint creation
from pptx.util import Inches  # For specifying dimensions
import os

def create_word_document(contents, output_file):
    doc = Document()
    for content in contents:
        if content['type'] == 'paragraph':
            doc.add_paragraph(content['text'])
        elif content['type'] == 'list':
            for item in content['items']:
                doc.add_paragraph(item, style='ListBullet')
        elif content['type'] == 'table':
            table = doc.add_table(rows=1, cols=len(content['header']))
            hdr_cells = table.rows[0].cells
            for i, header in enumerate(content['header']):
                hdr_cells[i].text = header
            for row_data in content['data']:
                row_cells = table.add_row().cells
                for i, cell_data in enumerate(row_data):
                    row_cells[i].text = str(cell_data)
    doc.save(output_file)

def create_powerpoint_presentation(contents, output_file):
    pres = Presentation()
    for content in contents:
        slide = pres.slides.add_slide(pres.slide_layouts[5])
        title = slide.shapes.title
        title.text = content['title']
        content_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
        content_frame = content_box.text_frame
        for bullet_point in content['bullet_points']:
            p = content_frame.add_paragraph()
            p.text = bullet_point
    pres.save(output_file)